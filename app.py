from flask import Flask, render_template, request, redirect, url_for, send_file
from werkzeug.utils import secure_filename
import os
from pdf2docx import Converter
from pdf2image import convert_from_path
from pptx import Presentation
from docx import Document
from openpyxl import Workbook, load_workbook
from PyPDF2 import PdfReader

# تحديث مسار Poppler
poppler_path = r"C:\Users\user\Downloads\poppler-24.08.0\Library\bin"

app = Flask(__name__)
app.config['UPLOAD_FOLDER'] = 'uploads'
app.config['CONVERTED_FOLDER'] = 'converted_files'

@app.route('/')
def home():
    """الصفحة الرئيسية."""
    return render_template('home.html')

@app.route('/convert/<filetype>', methods=['GET', 'POST'])
def convert(filetype):
    """تحويل الملفات بناءً على نوع التحويل."""
    if request.method == 'POST':
        file = request.files.get('file')
        if file:
            # حفظ الملف
            filename = secure_filename(file.filename)
            input_path = os.path.join(app.config['UPLOAD_FOLDER'], filename)
            file.save(input_path)

            try:
                # تنفيذ التحويل بناءً على نوع الملف
                if filetype == 'pdf_to_word':
                    output_filename = convert_pdf_to_word(input_path)
                elif filetype == 'pdf_to_excel':
                    output_filename = convert_pdf_to_excel(input_path)
                elif filetype == 'pdf_to_ppt':
                    output_filename = convert_pdf_to_ppt(input_path)
                elif filetype == 'to_pdf':
                    output_filename = convert_to_pdf(input_path)
                else:
                    return "Unsupported file type", 400

                output_path = os.path.join(app.config['CONVERTED_FOLDER'], output_filename)
                return send_file(output_path, as_attachment=True)
            except Exception as e:
                return f"Error: {str(e)}", 500
    return render_template('convert.html', filetype=filetype)

@app.route('/convert_pdf_to_word')
def route_convert_pdf_to_word():
    return redirect(url_for('convert', filetype='pdf_to_word'))

@app.route('/convert_pdf_to_excel')
def route_convert_pdf_to_excel():
    return redirect(url_for('convert', filetype='pdf_to_excel'))

@app.route('/convert_pdf_to_ppt')
def route_convert_pdf_to_ppt():
    return redirect(url_for('convert', filetype='pdf_to_ppt'))

@app.route('/convert_to_pdf')
def route_convert_to_pdf():
    return redirect(url_for('convert', filetype='to_pdf'))

def convert_pdf_to_word(input_path):
    """تحويل PDF إلى Word."""
    output_filename = os.path.splitext(os.path.basename(input_path))[0] + '.docx'
    output_path = os.path.join(app.config['CONVERTED_FOLDER'], output_filename)
    cv = Converter(input_path)
    cv.convert(output_path)
    cv.close()
    return output_filename

def convert_pdf_to_excel(input_path):
    """تحويل PDF إلى Excel."""
    output_filename = os.path.splitext(os.path.basename(input_path))[0] + '.xlsx'
    output_path = os.path.join(app.config['CONVERTED_FOLDER'], output_filename)

    wb = Workbook()
    ws = wb.active

    pdf = PdfReader(input_path)
    for page in pdf.pages:
        text = page.extract_text()
        for line in text.split('\n'):
            ws.append(line.split())

    wb.save(output_path)
    return output_filename

def convert_pdf_to_ppt(input_path):
    """تحويل PDF إلى PowerPoint."""
    output_filename = os.path.splitext(os.path.basename(input_path))[0] + '.pptx'
    images = convert_from_path(input_path, poppler_path=poppler_path)
    ppt = Presentation()
    for image in images:
        slide = ppt.slides.add_slide(ppt.slide_layouts[5])
        temp_image_path = os.path.join(app.config['UPLOAD_FOLDER'], "temp_image.jpg")
        image.save(temp_image_path)
        slide.shapes.add_picture(temp_image_path, 0, 0, ppt.slide_width, ppt.slide_height)
    output_path = os.path.join(app.config['CONVERTED_FOLDER'], output_filename)
    ppt.save(output_path)
    return output_filename

def convert_to_pdf(input_path):
    """تحويل ملفات أخرى إلى PDF."""
    output_filename = os.path.splitext(os.path.basename(input_path))[0] + '.pdf'
    output_path = os.path.join(app.config['CONVERTED_FOLDER'], output_filename)

    if input_path.lower().endswith('.docx'):
        doc = Document(input_path)
        doc.save(output_path)
    elif input_path.lower().endswith('.pptx'):
        ppt = Presentation(input_path)
        ppt.save(output_path)
    elif input_path.lower().endswith('.xlsx'):
        wb = load_workbook(input_path)
        ws = wb.active
        pdf_content = "\n".join(["\t".join([str(cell.value or "") for cell in row]) for row in ws.iter_rows()])
        with open(output_path, "w") as f:
            f.write(pdf_content)
    else:
        raise ValueError("Unsupported input file type")

@app.route('/register', methods=['GET', 'POST'])
def route_register():
    if request.method == 'POST':
        # هنا يمكنك إضافة الكود الخاص بمعالجة التسجيل مثل التحقق من البيانات، حفظ المستخدم في قاعدة البيانات
        return redirect(url_for('route_login'))  # إعادة توجيه بعد التسجيل
    return render_template('register.html')


if __name__ == '__main__':
    os.makedirs(app.config['UPLOAD_FOLDER'], exist_ok=True)
    os.makedirs(app.config['CONVERTED_FOLDER'], exist_ok=True)
    app.run(debug=True)
